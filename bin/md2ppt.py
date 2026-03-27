#!/usr/bin/env python3
"""
md2ppt.py — Markdown to PowerPoint converter

Markdown conventions:
  # Title        → title slide
  ## Heading     → new content slide
  ### Heading    → bold paragraph (level 0) within current slide
  - bullet       → bullet point
    - nested     → nested bullet (indented)
  [図解案] text  → collected into a dedicated "図解案" slide at the end

Usage:
  python bin/md2ppt.py md/sample.md ppt/output.pptx
"""

import sys
import re
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Error: python-pptx is required. Run: pip install python-pptx")
    sys.exit(1)


SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)

COLOR_TITLE_BG = RGBColor(0x1F, 0x49, 0x7D)   # 濃紺
COLOR_TITLE_FG = RGBColor(0xFF, 0xFF, 0xFF)   # 白
COLOR_HEADING_BG = RGBColor(0x2E, 0x75, 0xB6)  # 青
COLOR_HEADING_FG = RGBColor(0xFF, 0xFF, 0xFF)  # 白
COLOR_BODY_FG = RGBColor(0x1A, 0x1A, 0x1A)    # ほぼ黒
COLOR_ZU_BG = RGBColor(0xED, 0x7D, 0x31)      # オレンジ（図解案）


def add_title_slide(prs: Presentation, title: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_TITLE_BG

    txBox = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.5), Inches(10.33), Inches(2.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = COLOR_TITLE_FG


def add_content_slide(prs: Presentation, heading: str, body_items: list) -> None:
    """body_items: list of (indent_level: int, text: str, is_bold: bool)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ヘッダー帯
    hdr = slide.shapes.add_textbox(Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1))
    hdr.text_frame  # ensure initialized
    from pptx.util import Pt as _Pt
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = COLOR_HEADING_BG
    tf = hdr.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = f"  {heading}"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = COLOR_HEADING_FG

    if not body_items:
        return

    # 本文テキストボックス
    body_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.3), Inches(12.33), Inches(5.8)
    )
    tf = body_box.text_frame
    tf.word_wrap = True

    first = True
    for level, text, is_bold in body_items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.level = min(level, 4)
        indent_pt = Pt(12 + level * 16)
        p.space_before = Pt(2)

        run = p.add_run()
        if level == 0 and is_bold:
            run.text = text
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = COLOR_HEADING_BG
        else:
            bullet_char = "・" if level == 0 else "  ‒ " * (level)
            run.text = f"{bullet_char}{text}"
            run.font.size = Pt(14 if level == 0 else 12)
            run.font.bold = False
            run.font.color.rgb = COLOR_BODY_FG


def add_zuan_slide(prs: Presentation, zuan_items: list) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    hdr = slide.shapes.add_textbox(Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = COLOR_ZU_BG
    tf = hdr.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "  図解案"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = COLOR_TITLE_FG

    body_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.3), Inches(12.33), Inches(5.8)
    )
    tf = body_box.text_frame
    tf.word_wrap = True

    for i, (section, text) in enumerate(zuan_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"【{section}】 {text}"
        run.font.size = Pt(14)
        run.font.color.rgb = COLOR_BODY_FG


def parse_markdown(md_text: str):
    """
    Returns:
        slides: list of dict {type, title, body_items}
        zuan_items: list of (section_title, zuan_text)
    """
    slides = []
    zuan_items = []
    current_slide = None
    current_section = ""

    for raw_line in md_text.splitlines():
        line = raw_line.rstrip()

        # 図解案
        m = re.match(r'\[図解案\]\s*(.*)', line)
        if m:
            zuan_items.append((current_section, m.group(1)))
            continue

        # # Title slide
        m = re.match(r'^# (.+)', line)
        if m:
            current_section = m.group(1)
            slides.append({"type": "title", "title": m.group(1), "body_items": []})
            current_slide = slides[-1]
            continue

        # ## Content slide
        m = re.match(r'^## (.+)', line)
        if m:
            current_section = m.group(1)
            slides.append({"type": "content", "title": m.group(1), "body_items": []})
            current_slide = slides[-1]
            continue

        # ### Bold paragraph
        m = re.match(r'^### (.+)', line)
        if m and current_slide and current_slide["type"] == "content":
            current_slide["body_items"].append((0, m.group(1), True))
            continue

        # - bullet (nested by leading spaces, 2 spaces per level)
        m = re.match(r'^( *)[-*] (.+)', line)
        if m and current_slide and current_slide["type"] == "content":
            indent = len(m.group(1)) // 2 + 1
            current_slide["body_items"].append((indent, m.group(2), False))
            continue

    return slides, zuan_items


def convert(input_path: str, output_path: str) -> None:
    md_text = Path(input_path).read_text(encoding="utf-8")
    slides, zuan_items = parse_markdown(md_text)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for slide in slides:
        if slide["type"] == "title":
            add_title_slide(prs, slide["title"])
        elif slide["type"] == "content":
            add_content_slide(prs, slide["title"], slide["body_items"])

    if zuan_items:
        add_zuan_slide(prs, zuan_items)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"Saved: {output_path}  ({len(slides)} slides + {'図解案スライド' if zuan_items else '図解案なし'})")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python bin/md2ppt.py <input.md> <output.pptx>")
        sys.exit(1)
    convert(sys.argv[1], sys.argv[2])
